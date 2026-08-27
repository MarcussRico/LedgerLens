"""Read the saved .pptx back and emit an HTML rendering of it, so the real
artifact can be inspected and measured for text overflow."""
import html, base64, os
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400.0
SCALE = 96.0  # px per inch

FONTMAP = {"Segoe UI": "Inter, 'Helvetica Neue', Helvetica, Arial, sans-serif",
           "Georgia": "Georgia, 'Times New Roman', serif",
           "Consolas": "Menlo, Consolas, 'Courier New', monospace"}

def px(emu): return (emu / EMU_IN) * SCALE
def rgb(c):
    try: return "#%02X%02X%02X" % (c[0], c[1], c[2])
    except Exception: return None

prs = Presentation("LedgerLens_Deck.pptx")
SW, SH = px(prs.slide_width), px(prs.slide_height)
out = [f"""<!doctype html><meta charset=utf-8>
<style>
 body{{margin:0;background:#1a1a1a;font-family:Inter,sans-serif}}
 .slide{{position:relative;width:{SW}px;height:{SH}px;margin:24px auto;overflow:hidden;background:#0A0B0D}}
 .lbl{{color:#888;font:12px monospace;width:{SW}px;margin:6px auto 0}}
 .sh{{position:absolute;box-sizing:border-box}}
 .tx{{position:absolute;box-sizing:border-box;white-space:pre-wrap}}
 .ovf{{outline:2px solid #ff0066 !important}}
</style>"""]

for idx, slide in enumerate(prs.slides):
    out.append(f'<div class="lbl">slide {idx}</div><div class="slide" data-i="{idx}">')
    for shp in slide.shapes:
        L, T = px(shp.left or 0), px(shp.top or 0)
        Wd, Ht = px(shp.width or 0), px(shp.height or 0)
        if shp.shape_type == 13:  # picture
            img = shp.image.blob
            b64 = base64.b64encode(img).decode()
            out.append(f'<img class="sh" src="data:image/jpeg;base64,{b64}" '
                       f'style="left:{L}px;top:{T}px;width:{Wd}px;height:{Ht}px">')
            continue
        style = f"left:{L}px;top:{T}px;width:{Wd}px;height:{Ht}px;"
        # fill / line for autoshapes and connectors
        try:
            if shp.fill.type is not None and shp.fill.type == 1:
                style += f"background:{rgb(shp.fill.fore_color.rgb)};"
        except Exception: pass
        try:
            if shp.line.fill.type == 1:
                lw = (shp.line.width or 0) / 12700.0
                style += f"border:{max(lw,0.6)}px solid {rgb(shp.line.color.rgb)};"
        except Exception: pass
        has_text = shp.has_text_frame and shp.text_frame.text.strip()
        if not has_text:
            out.append(f'<div class="sh" style="{style}"></div>')
            continue
        inner = []
        for p in shp.text_frame.paragraphs:
            al = {1: "left", 2: "center", 3: "right"}.get(p.alignment, "left")
            ls = p.line_spacing or 1.2
            runs = []
            for child in p._p:
                if child.tag.endswith('}br'):
                    runs.append('<br>')
                    continue
                if not child.tag.endswith('}r'):
                    continue
                r = [x for x in p.runs if x._r is child]
                if not r: continue
                r = r[0]
                f = r.font
                fam = FONTMAP.get(f.name, f.name or "Inter, sans-serif")
                sz = (f.size.pt if f.size else 12) * 96.0 / 72.0   # pt -> css px
                col = rgb(f.color.rgb) if (f.color and f.color.type is not None) else "#E9E5DC"
                bold = "font-weight:600;" if f.bold else ""
                runs.append(f'<span style="font-family:{fam};font-size:{sz}px;color:{col};{bold}">'
                            f'{html.escape(r.text)}</span>')
            # line-height is relative to the element's own font-size, so the
            # paragraph must carry the run's size or every measurement is wrong
            base = 16.0
            for r0 in p.runs:
                if r0.font.size:
                    base = r0.font.size.pt * 96.0 / 72.0
                    break
            inner.append(f'<div style="text-align:{al};line-height:{ls};margin:0;font-size:{base}px">'
                         + "".join(runs) + "</div>")
        out.append(f'<div class="tx" style="{style}">' + "".join(inner) + "</div>")
    out.append("</div>")

out.append("""<script>
// flag any text box whose content is taller/wider than the box it was given
const bad = [];
document.querySelectorAll('.tx').forEach(el => {
  if (el.scrollHeight > el.clientHeight + 2 || el.scrollWidth > el.clientWidth + 2) {
    el.classList.add('ovf');
    bad.push({slide:+el.closest('.slide').dataset.i, text:el.innerText.slice(0,58).replace(/\\n/g,' | '),
              need:el.scrollHeight, got:el.clientHeight});
  }
});
window.__overflow = bad;
</script>""")
open("preview.html", "w").write("\n".join(out))
print("preview.html written")
